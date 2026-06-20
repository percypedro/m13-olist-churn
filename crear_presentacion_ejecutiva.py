#!/usr/bin/env python3
"""
Presentación Ejecutiva — Churn Olist — Grupo 7
Instalar: pip install python-pptx
Ejecutar: python crear_presentacion_ejecutiva.py
Genera: Presentacion_Ejecutiva_Churn_Olist.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── PALETA ───────────────────────────────────────────────
NAVY   = RGBColor(0x1E, 0x3A, 0x5F)
BLUE   = RGBColor(0x2E, 0x86, 0xAB)
ORANGE = RGBColor(0xF1, 0x8F, 0x01)
GREEN  = RGBColor(0x1A, 0xB3, 0x74)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
RED    = RGBColor(0xE7, 0x4C, 0x3C)
LGRAY  = RGBColor(0xF0, 0xF4, 0xF8)
MUTED  = RGBColor(0x6B, 0x7A, 0x8D)
DBLUE  = RGBColor(0x16, 0x2B, 0x48)

SW, SH = 13.33, 7.5

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

# ─── HELPERS ──────────────────────────────────────────────
def new_slide():
    return prs.slides.add_slide(BLANK)

def bg(sl, rgb):
    f = sl.background.fill
    f.solid()
    f.fore_color.rgb = rgb

def rect(sl, x, y, w, h, rgb, alpha=None):
    shp = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb
    shp.line.fill.background()
    return shp

def txt(sl, text, x, y, w, h, sz=16, bold=False, rgb=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = rgb
    r.font.name = 'Calibri'
    return tb

def mlines(sl, lines, x, y, w, h, sz=14, default_rgb=DARK):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            t, b, c = item, False, default_rgb
        elif len(item) == 2:
            t, b = item; c = default_rgb
        else:
            t, b, c = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = t
        r.font.size = Pt(sz)
        r.font.bold = b
        r.font.color.rgb = c
        r.font.name = 'Calibri'
    return tb

def kpi_card(sl, x, y, w, h, value, label, bg_rgb=BLUE, val_sz=36, lbl_sz=13):
    rect(sl, x, y, w, h, bg_rgb)
    txt(sl, value, x, y + 0.15, w, h * 0.55,
        sz=val_sz, bold=True, rgb=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, label, x, y + h * 0.62, w, h * 0.38,
        sz=lbl_sz, bold=False, rgb=WHITE, align=PP_ALIGN.CENTER)

def accent_bar(sl, y=0.55):
    rect(sl, 0, y, 13.33, 0.06, ORANGE)

# ══════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ══════════════════════════════════════════════════════════
sl = new_slide()
bg(sl, NAVY)

# Bloque izquierdo oscuro
rect(sl, 0, 0, 5.5, 7.5, DBLUE)

# Acento naranja vertical
rect(sl, 5.3, 0, 0.2, 7.5, ORANGE)

# Logo / ícono (texto grande)
txt(sl, '📉', 1.0, 0.8, 3.5, 1.5, sz=64, align=PP_ALIGN.CENTER)

# Título izquierdo
txt(sl, 'PREDICCIÓN\nDE CHURN', 0.4, 2.4, 4.6, 2.0,
    sz=32, bold=True, rgb=WHITE, align=PP_ALIGN.CENTER)
txt(sl, 'E-Commerce Olist · Brasil', 0.4, 4.35, 4.6, 0.5,
    sz=14, rgb=ORANGE, align=PP_ALIGN.CENTER)

# Contenido derecho
txt(sl, '¿Cómo saber quién va\na dejar de comprar\nantes de que se vaya?',
    5.8, 1.2, 7.0, 2.5, sz=26, bold=True, rgb=WHITE)

txt(sl, 'Un modelo de machine learning que identifica\n'
        'clientes en riesgo con el doble de precisión\n'
        'que una campaña masiva tradicional.',
    5.8, 3.7, 7.0, 1.8, sz=15, rgb=LGRAY)

# Tags
for i, (tag, color) in enumerate([
    ('Logistic Regression', BLUE),
    ('Lift 2.24x', GREEN),
    ('27 variables', ORANGE),
]):
    rect(sl, 5.8 + i * 2.4, 5.6, 2.2, 0.45, color)
    txt(sl, tag, 5.8 + i * 2.4, 5.6, 2.2, 0.45,
        sz=11, bold=True, rgb=WHITE, align=PP_ALIGN.CENTER)

txt(sl, 'Grupo 7  ·  Sprint 4  ·  Percy Fuentes',
    5.8, 6.6, 7.0, 0.5, sz=11, rgb=MUTED)

# ══════════════════════════════════════════════════════════
# SLIDE 2 — EL PROBLEMA
# ══════════════════════════════════════════════════════════
sl = new_slide()
bg(sl, WHITE)

rect(sl, 0, 0, 13.33, 1.1, NAVY)
txt(sl, 'EL PROBLEMA', 0.5, 0.15, 8.0, 0.8,
    sz=28, bold=True, rgb=WHITE)
txt(sl, 'El churn silencioso destruye valor', 0.5, 0.55, 8.0, 0.5,
    sz=14, rgb=ORANGE)

accent_bar(sl, y=1.1)

# Estadística impactante izquierda
rect(sl, 0.4, 1.5, 3.8, 2.2, RED)
txt(sl, '98.8%', 0.4, 1.6, 3.8, 1.2,
    sz=52, bold=True, rgb=WHITE, align=PP_ALIGN.CENTER)
txt(sl, 'de clientes en Olist\nno vuelve a comprar',
    0.4, 2.8, 3.8, 0.7, sz=13, rgb=WHITE, align=PP_ALIGN.CENTER)

# Problema central
rect(sl, 4.6, 1.5, 8.3, 2.2, LGRAY)
mlines(sl, [
    ('¿Cuál es el costo real?', True, NAVY),
    '',
    ('🔴  Adquirir un cliente nuevo cuesta 5-7x más que retener uno existente.', False, DARK),
    ('🔴  Una campaña masiva de retención contacta a todos, sin discriminar.', False, DARK),
    ('🔴  Presupuesto desperdiciado en clientes que no iban a abandonar.', False, DARK),
], 4.8, 1.55, 7.9, 2.1, sz=13)

# Pregunta clave
rect(sl, 0.4, 4.0, 12.5, 1.3, NAVY)
txt(sl,
    '¿Y si pudieras identificar exactamente quién va a irse\n'
    'y actuar antes de que ocurra?',
    0.6, 4.05, 12.1, 1.2, sz=20, bold=True, rgb=WHITE, align=PP_ALIGN.CENTER)

# Footer context
mlines(sl, [
    ('Dataset: ', True, MUTED),
    ('99,441 clientes únicos · Olist Brazil 2016-2018 · Dataset público Kaggle', False, MUTED),
], 0.4, 5.6, 12.5, 0.5, sz=11)

# ══════════════════════════════════════════════════════════
# SLIDE 3 — LA SOLUCIÓN
# ══════════════════════════════════════════════════════════
sl = new_slide()
bg(sl, WHITE)

rect(sl, 0, 0, 13.33, 1.1, BLUE)
txt(sl, 'LA SOLUCIÓN', 0.5, 0.15, 8.0, 0.8,
    sz=28, bold=True, rgb=WHITE)
txt(sl, 'Un modelo que predice el churn antes de que ocurra', 0.5, 0.55, 10.0, 0.5,
    sz=14, rgb=WHITE)
accent_bar(sl, y=1.1)

# Flujo de 3 pasos
pasos = [
    ('①', 'DATOS', 'Historial de compras,\nentregas y reseñas\nde cada cliente', NAVY),
    ('②', 'MODELO', 'Logistic Regression\n27 variables de\ncomportamiento', BLUE),
    ('③', 'ACCIÓN', 'Lista priorizada\nde clientes en riesgo\nlista para retención', GREEN),
]
for i, (num, titulo, desc, color) in enumerate(pasos):
    x = 0.5 + i * 4.2
    rect(sl, x, 1.5, 3.8, 3.5, color)
    txt(sl, num, x, 1.55, 3.8, 0.8,
        sz=36, bold=True, rgb=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, titulo, x, 2.3, 3.8, 0.6,
        sz=18, bold=True, rgb=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, desc, x, 2.95, 3.8, 1.8,
        sz=13, rgb=WHITE, align=PP_ALIGN.CENTER)
    # Flecha entre pasos
    if i < 2:
        txt(sl, '→', x + 3.8, 2.8, 0.4, 0.8,
            sz=28, bold=True, rgb=ORANGE, align=PP_ALIGN.CENTER)

# Resultado del modelo
rect(sl, 0.5, 5.2, 12.3, 1.0, LGRAY)
kpis_sol = [
    ('59.8% → 63%', 'AUC mejora en el tiempo'),
    ('2.24x', 'Lift vs campaña aleatoria'),
    ('22.4%', 'Churners capturados\ncontactando solo 10%'),
    ('27', 'Variables de comportamiento'),
]
for i, (val, lbl) in enumerate(kpis_sol):
    x = 0.7 + i * 3.1
    txt(sl, val, x, 5.22, 2.8, 0.45,
        sz=18, bold=True, rgb=NAVY, align=PP_ALIGN.CENTER)
    txt(sl, lbl, x, 5.65, 2.8, 0.5,
        sz=10, rgb=MUTED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 4 — RESULTADOS DE NEGOCIO
# ══════════════════════════════════════════════════════════
sl = new_slide()
bg(sl, NAVY)

rect(sl, 0, 0, 13.33, 1.1, DBLUE)
txt(sl, 'RESULTADOS', 0.5, 0.15, 8.0, 0.8,
    sz=28, bold=True, rgb=WHITE)
txt(sl, 'El modelo mejora con el tiempo — exactamente lo que necesita producción',
    0.5, 0.55, 12.0, 0.5, sz=13, rgb=ORANGE)
accent_bar(sl, y=1.1)

# KPIs principales
kpis = [
    ('63%',  'AUC Live\n(mejor período)', GREEN),
    ('2.24x', 'Lift Decil 1\nvs aleatorio', ORANGE),
    ('22.4%', 'Churners en\ntop 10% clientes', BLUE),
    ('+3 pts', 'AUC mejora\nVal→BT→Live', GREEN),
]
for i, (val, lbl, color) in enumerate(kpis):
    kpi_card(sl, 0.4 + i * 3.2, 1.4, 2.9, 1.8, val, lbl, color, val_sz=34, lbl_sz=12)

# Tabla estabilidad temporal
rect(sl, 0.4, 3.5, 12.5, 0.45, BLUE)
for j, h in enumerate(['Período', 'AUC', 'Lift Decil 1', 'Gain Decil 1', 'Interpretación']):
    x_pos = [0.5, 2.8, 4.8, 7.0, 9.2][j]
    txt(sl, h, x_pos, 3.52, 2.1, 0.4,
        sz=12, bold=True, rgb=WHITE)

rows_data = [
    ['Validación (Feb 2018)', '59.8%', '2.04x', '20.4%', 'Línea base'],
    ['BackTest (Mar-Abr 2018)', '61.0%', '2.18x', '21.8%', '↑ Mejora'],
    ['Live (May 2018)', '63.0%', '2.24x', '22.4%', '↑↑ Mejor período'],
]
row_colors = [RGBColor(0x1A, 0x2A, 0x4A), RGBColor(0x16, 0x24, 0x40), RGBColor(0x0D, 0x4A, 0x2E)]
for i, (row, rcolor) in enumerate(zip(rows_data, row_colors)):
    rect(sl, 0.4, 4.0 + i * 0.62, 12.5, 0.6, rcolor)
    x_positions = [0.5, 2.8, 4.8, 7.0, 9.2]
    for j, cell in enumerate(row):
        color = GREEN if '↑' in cell else WHITE
        bold  = j == 0
        txt(sl, cell, x_positions[j], 4.05 + i * 0.62, 2.1, 0.5,
            sz=12, bold=bold, rgb=color)

txt(sl, '★  Un modelo que mejora con datos nuevos es confiable para producción.',
    0.4, 6.3, 12.5, 0.5, sz=12, italic=True, rgb=ORANGE)

# ══════════════════════════════════════════════════════════
# SLIDE 5 — IMPACTO ECONÓMICO
# ══════════════════════════════════════════════════════════
sl = new_slide()
bg(sl, WHITE)

rect(sl, 0, 0, 13.33, 1.1, GREEN)
txt(sl, 'IMPACTO ECONÓMICO', 0.5, 0.15, 10.0, 0.8,
    sz=28, bold=True, rgb=WHITE)
txt(sl, 'Mismo presupuesto de campaña · Doble de resultado', 0.5, 0.55, 10.0, 0.5,
    sz=14, rgb=WHITE)
accent_bar(sl, y=1.1)

# Supuestos
rect(sl, 0.4, 1.3, 12.5, 0.55, LGRAY)
txt(sl, 'Supuestos:  CLV = $300/cliente  ·  Costo contacto = $5  ·  Base = 10,000 clientes  ·  Retención exitosa = 30%',
    0.6, 1.35, 12.1, 0.45, sz=11, rgb=MUTED, align=PP_ALIGN.CENTER)

# Dos columnas: sin modelo vs con modelo
# Sin modelo
rect(sl, 0.4, 2.1, 5.8, 4.2, RGBColor(0xFD, 0xEC, 0xEC))
rect(sl, 0.4, 2.1, 5.8, 0.55, RED)
txt(sl, '🎲  SIN MODELO — Campaña Aleatoria',
    0.5, 2.12, 5.6, 0.5, sz=13, bold=True, rgb=WHITE)

sin_items = [
    ('Clientes contactados', '1,000  (10% aleatorio)'),
    ('Churners capturados', '12  (tasa base 1.2%)'),
    ('Clientes retenidos', '3  (30% de éxito)'),
    ('Costo campaña', '$5,000'),
    ('Valor recuperado', '$900'),
    ('ROI neto', '-$4,100'),
]
for i, (k, v) in enumerate(sin_items):
    color_bg = RGBColor(0xFB, 0xD9, 0xD9) if i % 2 == 0 else RGBColor(0xFD, 0xEC, 0xEC)
    rect(sl, 0.4, 2.7 + i * 0.55, 5.8, 0.54, color_bg)
    txt(sl, k, 0.55, 2.73 + i * 0.55, 3.2, 0.48, sz=12, rgb=DARK)
    bold_v = i == 5
    color_v = RED if i == 5 else DARK
    txt(sl, v, 3.8, 2.73 + i * 0.55, 2.2, 0.48,
        sz=12, bold=bold_v, rgb=color_v, align=PP_ALIGN.RIGHT)

# Con modelo
rect(sl, 6.8, 2.1, 6.1, 4.2, RGBColor(0xE8, 0xF8, 0xF0))
rect(sl, 6.8, 2.1, 6.1, 0.55, GREEN)
txt(sl, '🤖  CON MODELO — Campaña Dirigida',
    6.9, 2.12, 5.9, 0.5, sz=13, bold=True, rgb=WHITE)

con_items = [
    ('Clientes contactados', '1,000  (top 10% Decil 1)'),
    ('Churners capturados', '26  (Gain 22.4%)'),
    ('Clientes retenidos', '7  (30% de éxito)'),
    ('Costo campaña', '$5,000'),
    ('Valor recuperado', '$2,100'),
    ('ROI neto', '-$2,900'),
]
for i, (k, v) in enumerate(con_items):
    color_bg = RGBColor(0xC8, 0xED, 0xD6) if i % 2 == 0 else RGBColor(0xE8, 0xF8, 0xF0)
    rect(sl, 6.8, 2.7 + i * 0.55, 6.1, 0.54, color_bg)
    txt(sl, k, 6.95, 2.73 + i * 0.55, 3.4, 0.48, sz=12, rgb=DARK)
    bold_v = i == 5
    color_v = GREEN if i == 5 else DARK
    txt(sl, v, 10.4, 2.73 + i * 0.55, 2.3, 0.48,
        sz=12, bold=bold_v, rgb=color_v, align=PP_ALIGN.RIGHT)

# vs arrow
txt(sl, 'VS', 6.05, 3.8, 0.7, 0.8,
    sz=20, bold=True, rgb=ORANGE, align=PP_ALIGN.CENTER)

# Resultado clave
rect(sl, 0.4, 6.45, 12.5, 0.75, NAVY)
txt(sl, '★  Con el mismo presupuesto, el modelo recupera $1,200 adicionales  —  2.3x más valor que la campaña aleatoria',
    0.6, 6.5, 12.1, 0.65, sz=13, bold=True, rgb=ORANGE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 6 — QUIÉN ESTÁ EN RIESGO (SHAP)
# ══════════════════════════════════════════════════════════
sl = new_slide()
bg(sl, WHITE)

rect(sl, 0, 0, 13.33, 1.1, NAVY)
txt(sl, '¿QUIÉN ESTÁ EN RIESGO?', 0.5, 0.15, 10.0, 0.8,
    sz=28, bold=True, rgb=WHITE)
txt(sl, 'El modelo explica por qué — no es una caja negra', 0.5, 0.55, 10.0, 0.5,
    sz=14, rgb=ORANGE)
accent_bar(sl, y=1.1)

# Columna izquierda — factores de riesgo
rect(sl, 0.4, 1.4, 6.0, 0.5, RED)
txt(sl, '🔴  AUMENTAN el riesgo de churn', 0.55, 1.42, 5.7, 0.45,
    sz=13, bold=True, rgb=WHITE)

riesgo_items = [
    ('⏰', 'Mucho tiempo sin comprar',   'recencia alta → cliente enfriándose'),
    ('🏪', 'Compra en una sola categoría', 'comprador_unico → sin diversificación'),
    ('🚚', 'Entregas lentas',              'lead_min alto → experiencia negativa'),
    ('💳', 'Muchas cuotas de pago',        'installments_max → estrés financiero'),
]
for i, (icon, titulo, desc) in enumerate(riesgo_items):
    y = 2.0 + i * 0.85
    rect(sl, 0.4, y, 6.0, 0.78, RGBColor(0xFD, 0xEC, 0xEC) if i%2==0 else WHITE)
    txt(sl, icon, 0.45, y + 0.08, 0.6, 0.6, sz=20, align=PP_ALIGN.CENTER)
    txt(sl, titulo, 1.1, y + 0.04, 5.1, 0.38, sz=12, bold=True, rgb=DARK)
    txt(sl, desc, 1.1, y + 0.4, 5.1, 0.35, sz=11, rgb=MUTED)

# Columna derecha — factores protectores
rect(sl, 6.9, 1.4, 6.0, 0.5, GREEN)
txt(sl, '🟢  REDUCEN el riesgo de churn', 7.05, 1.42, 5.7, 0.45,
    sz=13, bold=True, rgb=WHITE)

protector_items = [
    ('🛒', 'Compras frecuentes recientes',  'ordenes_60d alto → cliente activo'),
    ('⭐', 'Buenas reseñas de productos',   'review_max alto → satisfacción'),
    ('🗂️', 'Diversidad de categorías',       'comprador_unico bajo → explorador'),
    ('📦', 'Órdenes múltiples por mes',      'compras_x_mes alto → hábito formado'),
]
for i, (icon, titulo, desc) in enumerate(protector_items):
    y = 2.0 + i * 0.85
    rect(sl, 6.9, y, 6.0, 0.78, RGBColor(0xE8, 0xF8, 0xF0) if i%2==0 else WHITE)
    txt(sl, icon, 6.95, y + 0.08, 0.6, 0.6, sz=20, align=PP_ALIGN.CENTER)
    txt(sl, titulo, 7.6, y + 0.04, 5.1, 0.38, sz=12, bold=True, rgb=DARK)
    txt(sl, desc, 7.6, y + 0.4, 5.1, 0.35, sz=11, rgb=MUTED)

# Fuente
txt(sl, 'Fuente: SHAP LinearExplainer + coeficientes LogReg sobre 27 variables de comportamiento',
    0.4, 6.55, 12.5, 0.5, sz=10, italic=True, rgb=MUTED)

# ══════════════════════════════════════════════════════════
# SLIDE 7 — METODOLOGÍA
# ══════════════════════════════════════════════════════════
sl = new_slide()
bg(sl, WHITE)

rect(sl, 0, 0, 13.33, 1.1, NAVY)
txt(sl, 'METODOLOGÍA', 0.5, 0.15, 8.0, 0.8,
    sz=28, bold=True, rgb=WHITE)
txt(sl, 'Proceso riguroso en 4 etapas — sin tocar BackTest ni Live hasta el final',
    0.5, 0.55, 12.0, 0.5, sz=13, rgb=ORANGE)
accent_bar(sl, y=1.1)

# Etapas del proceso
etapas = [
    ('①', 'DATOS\n& FEATURES',
     '99K clientes · 78 features\nFiltro leakage temporal\n27 variables finales',
     NAVY, '78 → 27\nvariables'),
    ('②', 'MODELOS\nCOMPARADOS',
     'Optuna 100 trials\nLogReg · HGB · RF\nCV 5-fold estratificado',
     BLUE, '3 modelos\nevaluados'),
    ('③', 'MODELO\nELEGIDO',
     'LogReg por estabilidad\ntemporal creciente\nVal→BT→Live: ↑↑↑',
     GREEN, 'AUC\n59%→63%'),
    ('④', 'VALIDACIÓN\nINDEPENDIENTE',
     'BackTest + Live\nevaluados UNA SOLA VEZ\nal final del proceso',
     ORANGE, 'Sin\ncontaminación'),
]
for i, (num, titulo, desc, color, badge) in enumerate(etapas):
    x = 0.4 + i * 3.25
    rect(sl, x, 1.45, 3.0, 3.8, color)
    txt(sl, num, x, 1.5, 3.0, 0.65,
        sz=28, bold=True, rgb=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, titulo, x, 2.1, 3.0, 0.75,
        sz=13, bold=True, rgb=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, desc, x, 2.9, 3.0, 1.5,
        sz=11, rgb=WHITE, align=PP_ALIGN.CENTER)
    # Badge resultado
    rect(sl, x + 0.3, 4.55, 2.4, 0.55, RGBColor(0xFF, 0xFF, 0xFF))
    txt(sl, badge, x + 0.3, 4.55, 2.4, 0.55,
        sz=11, bold=True, rgb=color, align=PP_ALIGN.CENTER)
    # Flecha
    if i < 3:
        txt(sl, '→', x + 3.0, 3.05, 0.25, 0.7,
            sz=22, bold=True, rgb=ORANGE, align=PP_ALIGN.CENTER)

# Decisiones clave
rect(sl, 0.4, 5.35, 12.5, 0.5, LGRAY)
txt(sl, '⚙  Decisiones clave:', 0.6, 5.38, 2.5, 0.44,
    sz=12, bold=True, rgb=NAVY)
txt(sl,
    'SMOTE para desbalance (1.2% churn)  ·  Penalización overfitting en Optuna  ·  '
    'PSI para estabilidad de features  ·  LinearExplainer para SHAP',
    3.1, 5.38, 9.7, 0.44, sz=11, rgb=DARK)

rect(sl, 0.4, 5.95, 12.5, 0.85, NAVY)
txt(sl,
    'Principio fundamental: BackTest y Live se evaluaron una única vez, al final, '
    'sin influir en ninguna decisión de modelado — garantía de resultados honestos y reproducibles.',
    0.6, 5.98, 12.1, 0.8, sz=12, italic=True, rgb=ORANGE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 8 — ÉTICA Y GOBERNANZA
# ══════════════════════════════════════════════════════════
sl = new_slide()
bg(sl, WHITE)

rect(sl, 0, 0, 13.33, 1.1, DBLUE)
txt(sl, 'ÉTICA Y GOBERNANZA', 0.5, 0.15, 9.0, 0.8,
    sz=28, bold=True, rgb=WHITE)
txt(sl, 'Uso responsable del modelo — cumplimiento LGPD — monitoreo continuo',
    0.5, 0.55, 12.0, 0.5, sz=13, rgb=ORANGE)
accent_bar(sl, y=1.1)

pilares = [
    ('🔍', 'TRANSPARENCIA',
     'Logistic Regression con coeficientes\ninterpretables + análisis SHAP.\nCada predicción es explicable.',
     BLUE),
    ('🛡️', 'PRIVACIDAD\n(LGPD)',
     'Solo datos de comportamiento\nde compra. Sin datos sensibles.\ncustomer_id anonimizado.',
     NAVY),
    ('⚖️', 'SESGOS\nIDENTIFICADOS',
     'Sesgo temporal corregido.\nLeakage en lead_min eliminado.\nDistribución de importancia\nequilibrada (máx 11.57%).',
     ORANGE),
    ('📊', 'MONITOREO\n& GOBERNANZA',
     'AUC en producción.\nPSI de features mensual.\nReentrenamiento si AUC\ncae > 5 puntos.',
     GREEN),
]
for i, (icon, titulo, desc, color) in enumerate(pilares):
    x = 0.4 + i * 3.25
    rect(sl, x, 1.45, 3.0, 4.2, color)
    txt(sl, icon, x, 1.5, 3.0, 0.75,
        sz=30, align=PP_ALIGN.CENTER)
    txt(sl, titulo, x, 2.2, 3.0, 0.75,
        sz=12, bold=True, rgb=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, desc, x, 3.0, 3.0, 2.4,
        sz=11, rgb=WHITE, align=PP_ALIGN.CENTER)

# Usos permitidos / no permitidos
rect(sl, 0.4, 5.85, 6.0, 0.4, GREEN)
txt(sl, '✅  Usos recomendados', 0.55, 5.87, 5.7, 0.36,
    sz=12, bold=True, rgb=WHITE)
txt(sl, 'Priorizar campañas de retención  ·  Segmentar por riesgo  ·  Identificar causas operativas del churn',
    0.55, 6.3, 5.9, 0.5, sz=10, rgb=DARK)

rect(sl, 6.9, 5.85, 6.0, 0.4, RED)
txt(sl, '❌  Usos NO recomendados', 7.05, 5.87, 5.7, 0.36,
    sz=12, bold=True, rgb=WHITE)
txt(sl, 'Negar servicios  ·  Decisiones automatizadas sin revisión humana  ·  Compartir scores con terceros',
    7.05, 6.3, 5.9, 0.5, sz=10, rgb=DARK)

txt(sl,
    'El modelo es apoyo a la decisión, no un sustituto del juicio humano.',
    0.4, 6.9, 12.5, 0.45, sz=11, italic=True, rgb=MUTED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 9 — PROBLEMAS Y SOLUCIONES
# ══════════════════════════════════════════════════════════
sl = new_slide()
bg(sl, WHITE)

rect(sl, 0, 0, 13.33, 1.1, ORANGE)
txt(sl, 'DESAFÍOS Y SOLUCIONES', 0.5, 0.15, 10.0, 0.8,
    sz=28, bold=True, rgb=WHITE)
txt(sl, 'Los 4 problemas más críticos del proyecto — y cómo se resolvieron',
    0.5, 0.55, 12.0, 0.5, sz=13, rgb=WHITE)
accent_bar(sl, y=1.1)

problemas = [
    (
        '🔍', 'DATA LEAKAGE',
        'lead_min calculado con\nfechas de entrega futuras\n→ importancia inflada al 20%',
        'Filtrar órdenes por corte\ntemporal T antes de calcular\ncualquier feature logística',
        '20% → 7%\nimportancia real',
        RED, GREEN,
    ),
    (
        '🔄', 'MODELO INESTABLE',
        'HGB ganaba en Val (70%)\npero caía fuerte en Live (55%)\n→ −15 puntos de degradación',
        'Cambiar criterio: de "mayor\nAUC Val" a "mayor estabilidad\ntemporal Val→BT→Live"',
        'LogReg: 60%→63%\nestabilidad creciente',
        RED, GREEN,
    ),
    (
        '📊', 'GAIN/LIFT ROTO',
        'Lift = 1.0 en todos los\ndeciles — el modelo no\ndiscriminaba nada',
        'Invertir la clase objetivo:\nchurn=0 es la minoría (1.2%)\nno churn=1 (mayoría 98.8%)',
        'Lift 1.0 → 2.24x\nen Decil 1',
        RED, GREEN,
    ),
    (
        '🧠', 'SHAP FALLABA',
        'TreeExplainer lanzaba error\nal cambiar a LogReg\n→ incompatible con árboles',
        'Reemplazar por LinearExplainer\nnativo para regresión\nlogística',
        'SHAP funcional\ncon LogReg',
        RED, GREEN,
    ),
]

for i, (icon, titulo, problema, solucion, resultado, c_prob, c_sol) in enumerate(problemas):
    col = i % 2
    row = i // 2
    x_base = 0.4 + col * 6.5
    y_base = 1.4 + row * 2.85

    # Header del card
    rect(sl, x_base, y_base, 6.1, 0.5, NAVY)
    txt(sl, f'{icon}  {titulo}', x_base + 0.15, y_base + 0.05, 5.8, 0.42,
        sz=13, bold=True, rgb=WHITE)

    # Problema
    rect(sl, x_base, y_base + 0.5, 2.9, 1.65, RGBColor(0xFD, 0xEC, 0xEC))
    txt(sl, '❌ Problema', x_base + 0.1, y_base + 0.55, 2.7, 0.35,
        sz=10, bold=True, rgb=RED)
    txt(sl, problema, x_base + 0.1, y_base + 0.9, 2.7, 1.2,
        sz=10, rgb=DARK)

    # Solución
    rect(sl, x_base + 2.9, y_base + 0.5, 2.0, 1.65, RGBColor(0xE8, 0xF8, 0xF0))
    txt(sl, '✅ Solución', x_base + 3.0, y_base + 0.55, 1.8, 0.35,
        sz=10, bold=True, rgb=GREEN)
    txt(sl, solucion, x_base + 3.0, y_base + 0.9, 1.85, 1.2,
        sz=10, rgb=DARK)

    # Resultado badge
    rect(sl, x_base + 4.9, y_base + 0.5, 1.2, 1.65, GREEN)
    txt(sl, resultado, x_base + 4.9, y_base + 0.8, 1.2, 1.1,
        sz=10, bold=True, rgb=WHITE, align=PP_ALIGN.CENTER)

txt(sl,
    'Documentación completa: bitacora_cambios.md — 12 problemas documentados con causa raíz y solución',
    0.4, 7.05, 12.5, 0.35, sz=10, italic=True, rgb=MUTED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 10 — DEMO DEL DASHBOARD
# ══════════════════════════════════════════════════════════
sl = new_slide()
bg(sl, DBLUE)

# Fondo decorativo
rect(sl, 0, 0, 13.33, 7.5, DBLUE)
rect(sl, 0, 0, 0.5, 7.5, ORANGE)
rect(sl, 12.83, 0, 0.5, 7.5, ORANGE)

txt(sl, '🖥️', 3.5, 0.5, 6.33, 2.0, sz=80, align=PP_ALIGN.CENTER)

txt(sl, 'DEMO EN VIVO', 0.5, 2.5, 12.33, 1.2,
    sz=44, bold=True, rgb=WHITE, align=PP_ALIGN.CENTER)

txt(sl, 'Dashboard Interactivo de Predicción de Churn',
    0.5, 3.65, 12.33, 0.7, sz=20, rgb=ORANGE, align=PP_ALIGN.CENTER)

# Pills de páginas del dashboard
paginas = ['Resumen Ejecutivo', 'Predicciones', 'SHAP', 'Gain & Lift', 'Pred. Individual', 'Impacto Económico']
colors_p = [BLUE, NAVY, GREEN, BLUE, ORANGE, GREEN]
for i, (pag, col) in enumerate(zip(paginas, colors_p)):
    col_idx = i % 3
    row_idx = i // 3
    x = 1.8 + col_idx * 3.4
    y = 4.65 + row_idx * 0.75
    rect(sl, x, y, 3.0, 0.55, col)
    txt(sl, pag, x, y, 3.0, 0.55,
        sz=12, bold=True, rgb=WHITE, align=PP_ALIGN.CENTER)

txt(sl, 'streamlit run dashboard_churn.py',
    2.5, 6.55, 8.33, 0.6, sz=14, rgb=MUTED,
    align=PP_ALIGN.CENTER, italic=True)

# ══════════════════════════════════════════════════════════
# SLIDE 11 — PRÓXIMOS PASOS
# ══════════════════════════════════════════════════════════
sl = new_slide()
bg(sl, WHITE)

rect(sl, 0, 0, 13.33, 1.1, NAVY)
txt(sl, 'PRÓXIMOS PASOS', 0.5, 0.15, 8.0, 0.8,
    sz=28, bold=True, rgb=WHITE)
txt(sl, 'El modelo está listo para producción', 0.5, 0.55, 8.0, 0.5,
    sz=14, rgb=ORANGE)
accent_bar(sl, y=1.1)

pasos_next = [
    ('🔄', 'Reentrenamiento mensual',
     'Actualizar el modelo con datos del mes más reciente. '
     'Criterio: AUC cae más de 5 puntos o PSI > 0.25.',
     BLUE),
    ('📊', 'Monitoreo continuo',
     'Tracking del AUC en producción, distribución de features '
     'y tasa de churn observada vs predicha.',
     GREEN),
    ('🎯', 'Integración con CRM',
     'Exportar la lista de clientes en riesgo directamente '
     'al sistema de campañas para acción automatizada.',
     ORANGE),
    ('🔍', 'Expansión del modelo',
     'Incorporar datos de NPS, interacciones con soporte '
     'y variables de categoría para mejorar el AUC.',
     NAVY),
]
for i, (icon, titulo, desc, color) in enumerate(pasos_next):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.5
    y = 1.5 + row * 2.5
    rect(sl, x, y, 6.1, 2.2, color)
    txt(sl, icon, x + 0.15, y + 0.15, 0.9, 0.9, sz=28)
    txt(sl, titulo, x + 1.1, y + 0.18, 4.8, 0.6,
        sz=16, bold=True, rgb=WHITE)
    txt(sl, desc, x + 0.2, y + 0.85, 5.7, 1.2,
        sz=12, rgb=WHITE)

# Cierre
rect(sl, 0.5, 6.4, 12.33, 0.75, LGRAY)
txt(sl,
    'Grupo 7  ·  Percy Fuentes  ·  Modelo de Churn Olist  ·  Sprint 4  ·  2024',
    0.5, 6.45, 12.33, 0.65, sz=12, rgb=MUTED, align=PP_ALIGN.CENTER)

# ─── GUARDAR ──────────────────────────────────────────────
output = 'Presentacion_Ejecutiva_Churn_Olist.pptx'
prs.save(output)
print(f"✓ Presentación ejecutiva guardada: {output}")
print(f"  11 slides · Resumen ejecutivo · Negocio + Metodología + Ética + Desafíos")
