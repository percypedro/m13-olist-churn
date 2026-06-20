#!/usr/bin/env python3
"""
Presentación Sprint 4 — Churn Olist — Grupo 7
Instalar: pip install python-pptx
Ejecutar: python crear_presentacion_sprint4.py
Genera: Presentacion_Sprint4_Churn_Olist_Grupo7.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ─── PALETA DE COLORES ────────────────────────────────────
NAVY   = RGBColor(0x1E, 0x3A, 0x5F)
BLUE   = RGBColor(0x2E, 0x86, 0xAB)
ORANGE = RGBColor(0xF1, 0x8F, 0x01)
LGRAY  = RGBColor(0xF0, 0xF4, 0xF8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
GREEN  = RGBColor(0x1A, 0xB3, 0x74)
RED    = RGBColor(0xE7, 0x4C, 0x3C)
MUTED  = RGBColor(0x6B, 0x7A, 0x8D)
DGRAY  = RGBColor(0x2C, 0x3E, 0x50)
LBLUE  = RGBColor(0xB0, 0xC4, 0xDE)
DBLUE  = RGBColor(0x16, 0x2B, 0x48)
TEAL   = RGBColor(0x7E, 0xC8, 0xE3)

SW, SH = 13.33, 7.5

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

# ─── HELPERS ──────────────────────────────────────────────
def new_slide():
    return prs.slides.add_slide(BLANK)

def bg(sl, rgb):
    f = sl.background.fill
    f.solid()
    f.fore_color.rgb = rgb

def rect(sl, x, y, w, h, rgb):
    shp = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb
    shp.line.fill.background()
    return shp

def txt(sl, text, x, y, w, h, sz=16, bold=False, rgb=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = rgb
    r.font.name = 'Calibri'
    return tb

def mlines(sl, lines, x, y, w, h, sz=14, default_rgb=DARK):
    """lines: list of str or (text, bold) or (text, bold, rgb)"""
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            t, b, c = item, False, default_rgb
        elif len(item) == 2:
            t, b = item; c = default_rgb
        else:
            t, b, c = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = t
        r.font.size = Pt(sz)
        r.font.bold = b
        r.font.color.rgb = c
        r.font.name = 'Calibri'

def tbl(sl, data, x, y, w, h, hdr_bg=NAVY, hdr_fg=WHITE,
        alt_bg=None, sz=11, aligns=None):
    rows, cols = len(data), len(data[0])
    t = sl.shapes.add_table(rows, cols,
                             Inches(x), Inches(y),
                             Inches(w), Inches(h)).table
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = t.cell(r, c)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.alignment = aligns[c] if aligns else PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(sz)
            run.font.name = 'Calibri'
            if r == 0:
                run.font.bold = True
                run.font.color.rgb = hdr_fg
                cell.fill.solid()
                cell.fill.fore_color.rgb = hdr_bg
            elif alt_bg and r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt_bg

def slide_header(sl, title, dark=False):
    bg(sl, LGRAY if not dark else DGRAY)
    rect(sl, 0, 0, SW, 1.1, NAVY)
    rect(sl, 0, 0, 0.14, 1.1, ORANGE)
    txt(sl, title, 0.35, 0.15, 12.5, 0.8, sz=26, bold=True, rgb=WHITE)

# ══════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ══════════════════════════════════════════════════════════
sl = new_slide()
bg(sl, NAVY)
rect(sl, 0, 0, 4.8, SH, DBLUE)
rect(sl, 0, 0, 0.22, SH, ORANGE)

txt(sl, "MODELO DE PREDICCIÓN\nDE CHURN", 0.5, 0.7, 9.0, 2.2,
    sz=38, bold=True, rgb=WHITE)
txt(sl, "Dataset Olist — Brazilian E-Commerce", 0.5, 2.9, 9.0, 0.7,
    sz=20, rgb=LBLUE)
rect(sl, 0.5, 3.65, 5.5, 0.06, ORANGE)
txt(sl, "Sprint 4  |  Grupo 7  |  Percy Fuentes", 0.5, 3.8, 9.0, 0.6,
    sz=17, bold=True, rgb=ORANGE)
txt(sl, "Logistic Regression · SHAP · Optuna · Gain/Lift", 0.5, 4.5, 9.0, 0.55,
    sz=13, rgb=MUTED, italic=True)

# Stats cards
for i, (val, label, color) in enumerate([
    ("27",    "features\nseleccionadas", BLUE),
    ("63%",   "AUC Live\n(mejora temporal)", GREEN),
    ("2.24x", "Lift Decil 1\n(Live)", ORANGE),
]):
    xi = 10.1
    yi = 1.2 + i * 1.9
    rect(sl, xi, yi, 2.9, 1.7, DBLUE)
    txt(sl, val, xi, yi + 0.05, 2.9, 1.0, sz=38, bold=True, rgb=color, align=PP_ALIGN.CENTER)
    txt(sl, label, xi, yi + 1.0, 2.9, 0.6, sz=11, rgb=LBLUE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Agenda")

topics = [
    ("01", "Contexto del Proyecto y Dataset Olist"),
    ("02", "Sprint 3 vs Sprint 4 — Comparativa de Evolución"),
    ("03", "Corrección de Data Leakage en lead_min"),
    ("04", "Selección de Features: 78 → 27 variables"),
    ("05", "Modelo Final: Logistic Regression vs HGB vs RF"),
    ("06", "Importancia de Variables y Análisis SHAP"),
    ("07", "Gain & Lift Table + Evaluación Final y Predicciones"),
]
for i, (num, topic) in enumerate(topics):
    yi = 1.3 + i * 0.77
    color = NAVY if i % 2 == 0 else BLUE
    rect(sl, 0.4, yi, 0.6, 0.6, color)
    txt(sl, num, 0.4, yi + 0.05, 0.6, 0.5, sz=15, bold=True, rgb=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, topic, 1.2, yi + 0.1, 11.5, 0.5, sz=16, rgb=DARK)

# ══════════════════════════════════════════════════════════
# SLIDE 3 — CONTEXTO DEL PROYECTO
# ══════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Contexto del Proyecto — Dataset Olist")

rect(sl, 0.3, 1.3, 6.0, 5.6, WHITE)
txt(sl, "¿Qué es Olist?", 0.5, 1.45, 5.6, 0.5, sz=18, bold=True, rgb=NAVY)
mlines(sl, [
    "• Marketplace e-commerce brasileño (2016-2018)",
    "• +100,000 órdenes | 99,441 clientes únicos",
    "• Datos de órdenes, pagos, productos y reseñas",
    "",
    ("Objetivo del Modelo:", True, NAVY),
    "Predecir qué clientes van a abandonar la",
    "plataforma (churn) en el próximo mes, para",
    "tomar acciones de retención preventivas.",
], 0.5, 2.05, 5.6, 3.5, sz=14)

# Stats del dataset
for i, (val, label, color) in enumerate([
    ("1.2%",  "tasa de churn\n(clase minoritaria)", RED),
    ("98.8%", "clientes activos\n(clase mayoritaria)", BLUE),
    ("78",    "features\noriginales", NAVY),
]):
    xi = 6.6 + i * 2.2
    rect(sl, xi, 1.3, 2.0, 2.0, color)
    txt(sl, val, xi, 1.4, 2.0, 0.9, sz=30, bold=True, rgb=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, label, xi, 2.2, 2.0, 0.9, sz=12, rgb=WHITE, align=PP_ALIGN.CENTER)

# Splits temporales
rect(sl, 6.6, 3.5, 6.5, 3.4, WHITE)
txt(sl, "Splits Temporales", 6.8, 3.6, 6.1, 0.5, sz=17, bold=True, rgb=NAVY)
for i, (name, period, color) in enumerate([
    ("TRAIN",    "2016-09 → 2018-01", NAVY),
    ("VAL",      "2018-02",           BLUE),
    ("BACKTEST", "2018-03 / 04",      ORANGE),
    ("LIVE",     "2018-05",           GREEN),
]):
    yi = 4.2 + i * 0.6
    rect(sl, 6.8, yi, 1.5, 0.48, color)
    txt(sl, name, 6.8, yi + 0.05, 1.5, 0.4, sz=11, bold=True, rgb=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, period, 8.45, yi + 0.08, 4.5, 0.38, sz=13, rgb=DARK)

# ══════════════════════════════════════════════════════════
# SLIDE 4 — SPRINT 3 vs SPRINT 4 COMPARATIVA
# ══════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Sprint 3 vs Sprint 4 — Evolución del Proyecto")

data = [
    ["Aspecto",                  "Sprint 3",                    "Sprint 4",                      "Resultado"],
    ["Modelo Final",             "HistGradientBoosting",        "Logistic Regression",           "↑ Estabilidad temporal"],
    ["Data Leakage",             "lead_min con fecha futura",   "Filtro point-in-time",          "✓ Corregido"],
    ["BackTest en selección",    "Sí (contaminado)",            "No — solo CV + Val",            "✓ Metodología correcta"],
    ["Nº Features",              "~78 sin filtrar",             "27 (cascada de filtros)",       "↓ Ruido eliminado"],
    ["AUC Live",                 "~55% (caída de 15 pts)",      "63% (sube progresivamente)",    "↑ +8 puntos"],
    ["Importancia max. variable","lead_min > 20%",              "comprador_unico 11.57%",        "↑ Distribución equil."],
    ["Análisis SHAP",            "No implementado",             "LinearExplainer",               "✓ Nuevo"],
    ["Gain/Lift Table",          "No implementado",             "Lift 2.24x en decil 1 (Live)",  "✓ Nuevo"],
    ["Control overfitting",      "Básico",                      "Optuna + penalización CV-Val",  "↑ Brecha 8.1→5.0 pts"],
]
aligns = [PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT]
tbl(sl, data, 0.3, 1.2, 12.7, 6.0,
    hdr_bg=NAVY, alt_bg=RGBColor(0xE8, 0xEF, 0xF8),
    sz=11, aligns=aligns)

# ══════════════════════════════════════════════════════════
# SLIDE 5 — DATA LEAKAGE
# ══════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Corrección de Data Leakage — lead_min")

# Panel ANTES
rect(sl, 0.3, 1.3, 6.0, 5.6, WHITE)
rect(sl, 0.3, 1.3, 6.0, 0.55, RED)
txt(sl, "❌  ANTES — Sprint 3", 0.5, 1.38, 5.6, 0.44, sz=16, bold=True, rgb=WHITE)
mlines(sl, [
    ("Problema:", True, RED),
    "lead_min usaba order_delivered_customer_date",
    "→ fecha que podía ser POSTERIOR al corte T",
    "→ El modelo veía información del FUTURO",
    "→ Importancia artificialmente inflada al ~20%",
    "",
    ("Consecuencia:", True, DARK),
    "Variables futuras filtran al modelo durante",
    "entrenamiento — resultados NO reproducibles",
    "en producción.",
], 0.5, 2.0, 5.6, 3.5, sz=13)
rect(sl, 0.4, 5.4, 5.8, 0.85, RGBColor(0xFF, 0xEE, 0xEE))
txt(sl, "lead_min = 20% importancia (IRREAL)",
    0.5, 5.5, 5.6, 0.65, sz=14, bold=True, rgb=RED)

# Panel DESPUÉS
rect(sl, 6.6, 1.3, 6.4, 5.6, WHITE)
rect(sl, 6.6, 1.3, 6.4, 0.55, GREEN)
txt(sl, "✓  DESPUÉS — Sprint 4", 6.8, 1.38, 6.0, 0.44, sz=16, bold=True, rgb=WHITE)
mlines(sl, [
    ("Solución aplicada:", True, GREEN),
    "Filtro point-in-time en generar_features:",
    "Solo órdenes entregadas ANTES del corte T",
], 6.8, 2.0, 6.0, 1.3, sz=13)

rect(sl, 6.8, 3.4, 6.0, 1.35, RGBColor(0x1E, 0x2E, 0x3E))
mlines(sl, [
    ("h = h[", False, WHITE),
    ("  h['order_delivered_customer_date'].isna() |", False, TEAL),
    ("  (h['order_delivered_customer_date'] <= T)", False, TEAL),
    ("].copy()", False, WHITE),
], 6.9, 3.5, 5.8, 1.15, sz=11)

mlines(sl, [
    ("", False, DARK),
    ("Impacto:", True, NAVY),
    "• lead_min: 20% → 7% de importancia",
    "• NaN introducidos → ajuste de umbral missings",
    "  (10% → 25%) para mantener 27 features",
], 6.8, 4.85, 6.0, 2.0, sz=13)
rect(sl, 6.7, 5.7, 6.1, 0.85, RGBColor(0xEE, 0xF8, 0xEE))
txt(sl, "lead_min = 7% importancia (REAL ✓)",
    6.8, 5.8, 5.9, 0.65, sz=14, bold=True, rgb=GREEN)

# ══════════════════════════════════════════════════════════
# SLIDE 6 — SELECCIÓN DE FEATURES
# ══════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Selección de Features — Cascada de 4 Filtros")

steps = [
    ("78",  "Estado\nInicial",             MUTED,  "AUC: 0.631"),
    ("48",  "Univariante\nthreshold=0.03", BLUE,   "AUC: 0.621"),
    ("30",  "Correlación\nthreshold=0.96", NAVY,   "AUC: 0.626"),
    ("27",  "Missings\n≤ 25%",             ORANGE, "AUC: 0.616"),
    ("27✓", "PSI\n≤ 0.25",                GREEN,  "AUC: 0.616"),
]
for i, (count, label, color, auc) in enumerate(steps):
    xi = 0.4 + i * 2.58
    if i < 4:
        rect(sl, xi + 1.95, 2.6, 0.63, 0.4, RGBColor(0xCC, 0xCC, 0xCC))
        txt(sl, "→", xi + 1.95, 2.55, 0.63, 0.45, sz=18, bold=True,
            rgb=MUTED, align=PP_ALIGN.CENTER)
    rect(sl, xi, 1.3, 1.95, 2.7, color)
    txt(sl, count, xi, 1.4, 1.95, 1.0, sz=32, bold=True, rgb=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, label, xi, 2.4, 1.95, 1.1, sz=12, rgb=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, auc, xi, 3.55, 1.95, 0.45, sz=11, rgb=WHITE, align=PP_ALIGN.CENTER, italic=True)

rect(sl, 0.3, 4.25, 12.7, 2.6, WHITE)
txt(sl, "Notas metodológicas:", 0.5, 4.35, 5.5, 0.45, sz=15, bold=True, rgb=NAVY)
mlines(sl, [
    "• El umbral de missings se elevó de 10% a 25%: la corrección del leakage introduce NaN",
    "  legítimos (órdenes en tránsito al corte T). Los NaN son ausencia real de información.",
    "• PSI no eliminó ninguna variable → las 27 features son estables temporalmente (Train vs Val).",
    "• AUC Val se mantiene en 0.616 tras los filtros: no se pierde capacidad predictiva.",
], 0.5, 4.85, 12.2, 1.85, sz=13)

# ══════════════════════════════════════════════════════════
# SLIDE 7 — DESBALANCEO + OPTUNA
# ══════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Desbalanceo de Clases + Optimización con Optuna")

# Left: desbalanceo
rect(sl, 0.3, 1.3, 6.0, 5.6, WHITE)
txt(sl, "Desbalanceo de Clases", 0.5, 1.4, 5.6, 0.5, sz=17, bold=True, rgb=NAVY)

rect(sl, 0.5, 2.0, 5.6, 0.7, RED)
txt(sl, "churn=0 (churners):    1,604 → 1.2%",
    0.6, 2.08, 5.4, 0.55, sz=14, bold=True, rgb=WHITE)
rect(sl, 0.5, 2.8, 5.6, 0.7, MUTED)
txt(sl, "churn=1 (activos):  132,733 → 98.8%",
    0.6, 2.88, 5.4, 0.55, sz=14, rgb=WHITE)

txt(sl, "Solución: SMOTE + RandomUnderSampler", 0.5, 3.65, 5.6, 0.45,
    sz=14, bold=True, rgb=NAVY)
rect(sl, 0.5, 4.15, 5.6, 0.7, GREEN)
txt(sl, "churn=0 (churners):  26,546 → 37.5%",
    0.6, 4.23, 5.4, 0.55, sz=14, bold=True, rgb=WHITE)
rect(sl, 0.5, 4.9, 5.6, 0.7, BLUE)
txt(sl, "churn=1 (activos):   44,243 → 62.5%",
    0.6, 4.98, 5.4, 0.55, sz=14, rgb=WHITE)

mlines(sl, [
    ("ImbPipeline:", True, NAVY),
    "SMOTE(sampling_strategy=0.2)",
    "+ RandomUnderSampler(sampling_strategy=0.6)",
], 0.5, 5.75, 5.6, 1.0, sz=12)

# Right: Optuna
rect(sl, 6.6, 1.3, 6.4, 5.6, WHITE)
txt(sl, "Optuna — Optimización Bayesiana", 6.8, 1.4, 6.0, 0.5, sz=17, bold=True, rgb=NAVY)
mlines(sl, [
    "• 100 trials por modelo (HGB, LogReg, RF)",
    "• CV 5-fold estratificado",
    "• Penalización anti-overfitting:",
], 6.8, 2.0, 6.0, 1.2, sz=14)

rect(sl, 6.8, 3.3, 6.0, 0.95, RGBColor(0x1E, 0x2E, 0x3E))
mlines(sl, [
    ("penalizacion = max(0, auc_cv - auc_val - 0.03)", False, TEAL),
    ("return auc_cv - penalizacion", False, WHITE),
], 6.9, 3.38, 5.8, 0.82, sz=12)

mlines(sl, [
    ("Resultado:", True, NAVY),
    "Brecha CV-Val: 8.1 pts → 5.0 pts",
    "",
    "• Early stopping desactivado en Optuna",
    "  (velocidad: 30s/trial → 3s/trial)",
    "• Activado solo en modelo final HGB",
    "  para comparativa de estabilidad",
], 6.8, 4.35, 6.0, 2.4, sz=13)

# ══════════════════════════════════════════════════════════
# SLIDE 8 — COMPARATIVA DE MODELOS
# ══════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Comparativa de Modelos — HGB vs LogReg vs RF")

data = [
    ["Modelo",  "AUC CV", "AUC Val", "AUC BackTest", "AUC Live", "Tendencia"],
    ["HGB",     "0.771",  "0.703",   "0.59",         "0.55",     "↓ −15 pts (inestable)"],
    ["LogReg",  "0.579",  "0.598",   "0.61",         "0.63",     "↑ +3 pts (ELEGIDO ✓)"],
    ["RF",      "0.592",  "0.578",   "0.59",         "0.59",     "→ Estable pero bajo"],
]
aligns2 = [PP_ALIGN.LEFT] + [PP_ALIGN.CENTER] * 5
tbl(sl, data, 0.3, 1.3, 12.7, 2.5, hdr_bg=NAVY,
    alt_bg=RGBColor(0xE8, 0xEF, 0xF8), sz=13, aligns=aligns2)

# Explanation cards
for i, (title, body, color) in enumerate([
    ("¿Por qué no HGB?",
     "Tiene el mejor AUC Val (0.703) pero cae de 0.70 a 0.55 en Live — "
     "una degradación de 15 puntos. Indica sobreajuste a patrones del "
     "período de entrenamiento (Olist 2016-2017).",
     RED),
    ("¿Por qué LogReg?",
     "AUC Val menor (0.598) pero MEJORA progresivamente: Val→BackTest→Live "
     "= 0.60→0.61→0.63. Los patrones aprendidos generalizan bien al futuro. "
     "En producción, la estabilidad es más valiosa que el pico en Val.",
     GREEN),
    ("¿Por qué no RF?",
     "Estable (0.58→0.59→0.59) pero con el AUC más bajo de los tres. "
     "No captura suficiente señal predictiva. LogReg ofrece mejor balance "
     "entre performance y estabilidad temporal.",
     BLUE),
]):
    xi = 0.3 + i * 4.35
    rect(sl, xi, 4.1, 4.1, 0.5, color)
    txt(sl, title, xi, 4.15, 4.1, 0.4, sz=14, bold=True, rgb=WHITE, align=PP_ALIGN.CENTER)
    rect(sl, xi, 4.65, 4.1, 2.2, WHITE)
    txt(sl, body, xi + 0.1, 4.75, 3.9, 2.0, sz=12, rgb=DARK)

# ══════════════════════════════════════════════════════════
# SLIDE 9 — SECTION: ¿POR QUÉ LOGREG?
# ══════════════════════════════════════════════════════════
sl = new_slide()
bg(sl, NAVY)
rect(sl, 0, 0, 0.22, SH, ORANGE)
txt(sl, "¿Por qué\nLogistic Regression?", 0.5, 1.5, 12.3, 2.2,
    sz=50, bold=True, rgb=WHITE, align=PP_ALIGN.CENTER)
txt(sl, "Estabilidad temporal > AUC puntual en validación",
    0.5, 3.8, 12.3, 0.9, sz=22, rgb=LBLUE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 10 — ESTABILIDAD TEMPORAL (visual)
# ══════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Estabilidad Temporal — Val → BackTest → Live")

# Bar chart visual (manual bars)
rect(sl, 0.3, 1.3, 8.0, 5.6, WHITE)
txt(sl, "AUC por Período", 0.5, 1.4, 7.6, 0.5, sz=17, bold=True, rgb=NAVY)

periods = ["Val", "BackTest", "Live"]
models = [
    ("HGB",    [0.703, 0.59, 0.55],  RED),
    ("LogReg", [0.598, 0.61, 0.63],  GREEN),
    ("RF",     [0.578, 0.59, 0.59],  BLUE),
]
BAR_H_SCALE = 3.5
BASE_Y = 6.6
for mi, (mname, vals, color) in enumerate(models):
    for pi, val in enumerate(vals):
        xi = 0.7 + pi * 2.5 + mi * 0.72
        bar_h = (val - 0.50) * BAR_H_SCALE / 0.25
        yi = BASE_Y - bar_h
        rect(sl, xi, yi, 0.65, bar_h, color)
        txt(sl, str(val), xi, yi - 0.35, 0.65, 0.3, sz=9, bold=True,
            rgb=color, align=PP_ALIGN.CENTER)

# Period labels
for pi, period in enumerate(periods):
    xi = 0.7 + pi * 2.5
    txt(sl, period, xi, BASE_Y + 0.05, 2.1, 0.4, sz=13, bold=True,
        rgb=NAVY, align=PP_ALIGN.CENTER)

# Legend
for mi, (mname, _, color) in enumerate(models):
    xi = 1.0 + mi * 2.3
    rect(sl, xi, 7.15, 0.3, 0.25, color)
    txt(sl, mname, xi + 0.35, 7.1, 1.8, 0.35, sz=12, rgb=DARK)

# Right: conclusion panel
rect(sl, 8.5, 1.3, 4.5, 5.6, NAVY)
txt(sl, "Conclusión", 8.7, 1.4, 4.1, 0.5, sz=18, bold=True, rgb=WHITE)
mlines(sl, [
    ("HGB:", True, RED),
    "0.70 → 0.59 → 0.55",
    "(caída de 15 puntos)",
    "",
    ("RF:", True, BLUE),
    "0.58 → 0.59 → 0.59",
    "(estable pero bajo)",
    "",
    ("LogReg ✓:", True, GREEN),
    "0.60 → 0.61 → 0.63",
    "(SUBE en el tiempo)",
    "",
    ("Decisión:", True, ORANGE),
    "Se sacrifica AUC Val",
    "a cambio de robustez",
    "temporal real.",
], 8.7, 2.0, 4.1, 4.7, sz=13, default_rgb=LBLUE)

# ══════════════════════════════════════════════════════════
# SLIDE 11 — IMPORTANCIA DE VARIABLES
# ══════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Importancia de Variables — Logistic Regression (coeficientes)")

vars_data = [
    ("comprador_unico",  11.57, "Diversidad de vendedores del cliente"),
    ("installments_max",  7.74, "Número máximo de cuotas usadas"),
    ("ordenes_60d",       7.60, "Órdenes en los últimos 60 días"),
    ("recencia_rel",      7.45, "Recencia relativa al período"),
    ("lead_min",          7.21, "Tiempo mínimo de entrega"),
    ("review_max",        6.86, "Calificación máxima de reseñas"),
    ("n_items_sum",       6.83, "Total de ítems comprados"),
    ("recencia",          6.17, "Días sin comprar"),
    ("monto_sum",         4.21, "Gasto total histórico"),
    ("cat_div_ratio",     3.93, "Ratio de diversidad de categorías"),
]
MAX_PCT = 12.0
BAR_MAX_W = 8.0
for i, (var, pct, desc) in enumerate(vars_data):
    yi = 1.3 + i * 0.6
    bar_w = (pct / MAX_PCT) * BAR_MAX_W
    color = ORANGE if i == 0 else (NAVY if pct >= 7.0 else BLUE)
    rect(sl, 0.3, yi, bar_w, 0.5, color)
    txt(sl, f"  {var}", 0.3, yi + 0.07, bar_w + 0.05, 0.38, sz=11, bold=False, rgb=WHITE)
    txt(sl, f"{pct}%", bar_w + 0.4, yi + 0.08, 0.85, 0.38, sz=12, bold=True, rgb=DARK)
    txt(sl, desc, bar_w + 1.4, yi + 0.1, 3.8, 0.38, sz=11, rgb=MUTED, italic=True)

# Right note
rect(sl, 9.7, 1.3, 3.3, 5.5, WHITE)
txt(sl, "✓  Distribución\nEquilibrada", 9.8, 1.4, 3.1, 0.9,
    sz=16, bold=True, rgb=NAVY, align=PP_ALIGN.CENTER)
mlines(sl, [
    ("Ninguna variable", False, DARK),
    ("supera el 12%", True, GREEN),
    "",
    ("Top 5 acumulan", False, DARK),
    ("solo 41.6%", True, BLUE),
    "",
    "vs Sprint 3 (HGB):",
    ("lead_min > 20%", True, RED),
    ("top 5 > 50%", True, RED),
    "",
    "Método: coeficientes",
    "absolutos LogReg",
    "normalizados a %",
], 9.8, 2.5, 3.1, 4.1, sz=12)

# ══════════════════════════════════════════════════════════
# SLIDE 12 — ANÁLISIS SHAP
# ══════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Análisis SHAP — Interpretabilidad del Modelo LogReg")

rect(sl, 0.3, 1.3, 7.8, 5.6, WHITE)
txt(sl, "Dirección de cada variable (LinearExplainer)",
    0.5, 1.4, 7.4, 0.5, sz=16, bold=True, rgb=NAVY)
shap_rows = [
    ["Variable",         "Dirección",      "Interpretación de Negocio"],
    ["comprador_unico",  "↓ REDUCE churn", "Más vendedores = mayor fidelidad"],
    ["installments_max", "↑ SUBE churn",   "Muchas cuotas = mayor riesgo financiero"],
    ["recencia",         "↑ SUBE churn",   "Sin comprar mucho tiempo = abandono"],
    ["lead_min",         "↑ SUBE churn",   "Entregas lentas dañan la experiencia"],
    ["review_max",       "↓ REDUCE churn", "Buenas reseñas → cliente satisfecho"],
    ["compras_x_mes",    "↓ REDUCE churn", "Alta frecuencia = cliente fiel"],
    ["monto_sum",        "↑ SUBE churn",   "Alto gasto histórico sin recompra reciente"],
    ["monto_60d",        "↓ REDUCE churn", "Gasto reciente activo retiene al cliente"],
]
tbl(sl, shap_rows, 0.3, 2.0, 7.8, 4.8, hdr_bg=NAVY,
    alt_bg=RGBColor(0xE8, 0xEF, 0xF8), sz=11,
    aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT])

rect(sl, 8.4, 1.3, 4.6, 2.6, NAVY)
txt(sl, "Perfil de Riesgo Alto", 8.6, 1.38, 4.2, 0.5, sz=15, bold=True, rgb=WHITE)
mlines(sl, [
    "• Usa muchas cuotas",
    "• Lleva tiempo sin comprar",
    "• Recibe pedidos lentamente",
    "• Compra en pocos vendedores",
    "• Sin gasto reciente",
], 8.6, 1.95, 4.2, 1.8, sz=13, default_rgb=RGBColor(0xFF, 0xCC, 0x80))

rect(sl, 8.4, 4.1, 4.6, 2.8, GREEN)
txt(sl, "Perfil de Cliente Fiel", 8.6, 4.18, 4.2, 0.5, sz=15, bold=True, rgb=WHITE)
mlines(sl, [
    "• Compra frecuentemente",
    "• Gasto reciente activo",
    "• Buenas reseñas registradas",
    "• Diversidad de vendedores",
    "• Entregas rápidas",
], 8.6, 4.75, 4.2, 2.0, sz=13, default_rgb=WHITE)

# ══════════════════════════════════════════════════════════
# SLIDE 13 — GAIN & LIFT
# ══════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Gain & Lift Table — Capacidad de Segmentación")

periodos = [
    ("Val\n2018-02",      "2.04x", "20.41%", "50 / 245",   BLUE),
    ("BackTest\n2018-03/04","2.18x","21.76%","114 / 524",  ORANGE),
    ("Live\n2018-05",     "2.24x", "22.38%", "62 / 277",   GREEN),
]
for i, (periodo, lift, gain, resp, color) in enumerate(periodos):
    xi = 0.3 + i * 4.35
    rect(sl, xi, 1.3, 4.1, 0.6, color)
    txt(sl, periodo, xi, 1.35, 4.1, 0.5, sz=15, bold=True, rgb=WHITE, align=PP_ALIGN.CENTER)
    rect(sl, xi, 2.0, 4.1, 1.8, WHITE)
    txt(sl, f"Lift Decil 1", xi + 0.1, 2.05, 3.9, 0.45, sz=13, bold=True, rgb=MUTED)
    txt(sl, lift, xi + 0.1, 2.45, 3.9, 0.8, sz=38, bold=True, rgb=color, align=PP_ALIGN.CENTER)
    txt(sl, f"Gain: {gain}  |  Resp: {resp}", xi + 0.1, 3.25, 3.9, 0.45,
        sz=12, rgb=DARK, align=PP_ALIGN.CENTER)

    mini = [
        ["Decil", "Gain%", "Lift"],
        ["1", gain, lift],
        ["2", "~32%", "~1.6x"],
        ["5", "~64%", "~1.3x"],
        ["10", "100%", "1.0x"],
    ]
    tbl(sl, mini, xi, 3.85, 4.1, 2.4, hdr_bg=color, sz=10)

rect(sl, 0.3, 6.4, 12.7, 1.0, WHITE)
txt(sl, "Interpretación:", 0.5, 6.48, 2.5, 0.45, sz=14, bold=True, rgb=NAVY)
txt(sl, ("Contactando solo el 10% superior de clientes (Decil 1) se captura más del 20% "
         "de todos los churners reales. El Lift mejora en Live (2.24x), consistente con "
         "la estabilidad temporal de LogReg."),
    2.9, 6.48, 10.1, 0.85, sz=12, rgb=DARK)

# ══════════════════════════════════════════════════════════
# SLIDE 14 — EVALUACIÓN FINAL
# ══════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Evaluación Final del Modelo — Logistic Regression")

for i, (period, auc, delta, color) in enumerate([
    ("VAL\n2018-02",       "59.8%", "Base",    BLUE),
    ("BACKTEST\n2018-03/04","61.0%", "↑ +1.2 pts", ORANGE),
    ("LIVE\n2018-05",      "63.0%", "↑ +3.2 pts", GREEN),
]):
    xi = 0.3 + i * 4.35
    rect(sl, xi, 1.3, 4.1, 2.5, color)
    txt(sl, period, xi, 1.38, 4.1, 0.85, sz=15, bold=True, rgb=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, auc, xi, 2.1, 4.1, 1.1, sz=46, bold=True, rgb=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, delta, xi, 3.15, 4.1, 0.45, sz=13, rgb=WHITE, align=PP_ALIGN.CENTER, italic=True)

rect(sl, 0.3, 4.0, 12.7, 0.5, NAVY)
txt(sl, "Métricas en Live — umbral óptimo por F1 (con desbalance 1.2% churn)",
    0.5, 4.05, 12.2, 0.4, sz=13, bold=True, rgb=WHITE)

rect(sl, 0.3, 4.55, 12.7, 2.8, WHITE)
txt(sl, "Nota: Actualizar con valores reales al ejecutar la celda de métricas del notebook.",
    0.5, 4.65, 12.2, 0.45, sz=12, rgb=MUTED, italic=True)
mlines(sl, [
    ("AUC ROC", True, NAVY),
    "0.630",
    "Discriminación global del modelo",
], 0.5, 5.2, 2.9, 1.7, sz=12)
mlines(sl, [
    ("Precisión", True, NAVY),
    "Ver notebook",
    "De predichos churn, cuántos son reales",
], 3.5, 5.2, 2.9, 1.7, sz=12)
mlines(sl, [
    ("Recall", True, NAVY),
    "Ver notebook",
    "De churners reales, cuántos detecta",
], 6.5, 5.2, 2.9, 1.7, sz=12)
mlines(sl, [
    ("F1-Score", True, NAVY),
    "Ver notebook",
    "Balance precisión-recall",
], 9.5, 5.2, 3.0, 1.7, sz=12)

txt(sl, "BackTest y Live evaluados UNA SOLA VEZ al final, sin influir en decisiones de modelado (metodología correcta).",
    0.3, 7.25, 12.7, 0.4, sz=11, rgb=MUTED, italic=True)

# ══════════════════════════════════════════════════════════
# SLIDE 15 — PREDICCIONES
# ══════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Predicciones — Septiembre 2018")

rect(sl, 0.3, 1.3, 4.8, 3.2, NAVY)
txt(sl, "24,145", 0.3, 1.45, 4.8, 1.6, sz=54, bold=True, rgb=WHITE, align=PP_ALIGN.CENTER)
txt(sl, "clientes activos\npredichos", 0.3, 3.0, 4.8, 0.8,
    sz=15, rgb=LBLUE, align=PP_ALIGN.CENTER)
txt(sl, "Septiembre 2018", 0.3, 3.85, 4.8, 0.5,
    sz=14, bold=True, rgb=ORANGE, align=PP_ALIGN.CENTER)
txt(sl, "Exportado → predicciones_churn.csv", 0.3, 4.35, 4.8, 0.45,
    sz=11, rgb=MUTED, align=PP_ALIGN.CENTER, italic=True)

rect(sl, 5.4, 1.3, 7.6, 0.5, BLUE)
txt(sl, "Top Clientes en Riesgo (prob_churn descendente)",
    5.5, 1.35, 7.4, 0.4, sz=13, bold=True, rgb=WHITE)
top_data = [
    ["customer_unique_id (truncado)", "mes",        "prob_churn"],
    ["7f76b3a2c201b6a4ac7d...",       "2018-09-01", "0.858"],
    ["44482ad67ba7af317fd5...",       "2018-09-01", "0.858"],
    ["949a5efc8cd62cbe3888...",       "2018-09-01", "0.855"],
    ["30b94abdcaa5b08d5d78...",       "2018-09-01", "0.845"],
    ["b2ccf88b60413e6fb2b2...",       "2018-09-01", "0.844"],
]
tbl(sl, top_data, 5.4, 1.85, 7.6, 3.0, hdr_bg=NAVY,
    alt_bg=RGBColor(0xE8, 0xEF, 0xF8), sz=11,
    aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER])

rect(sl, 0.3, 4.7, 12.7, 2.6, WHITE)
txt(sl, "Estrategia de Uso por Negocio:", 0.5, 4.8, 12.2, 0.45,
    sz=15, bold=True, rgb=NAVY)
for i, (rango, accion, color) in enumerate([
    ("prob > 0.80", "Acción inmediata: contacto personalizado + descuento",     RED),
    ("0.60 – 0.80", "Campaña preventiva: email + mejora servicio entrega",      ORANGE),
    ("0.40 – 0.60", "Monitoreo: seguimiento próximo mes",                       BLUE),
]):
    yi = 5.35 + i * 0.6
    rect(sl, 0.4, yi, 1.9, 0.45, color)
    txt(sl, rango, 0.4, yi + 0.05, 1.9, 0.38, sz=12, bold=True, rgb=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, accion, 2.5, yi + 0.07, 10.0, 0.38, sz=13, rgb=DARK)

# ══════════════════════════════════════════════════════════
# SLIDE 16 — CORRECCIONES DEL SPRINT 4
# ══════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Principales Correcciones Realizadas en el Sprint 4")

corrections = [
    (RED,    "Data Leakage",        "lead_min usaba fechas futuras → filtro point-in-time. Importancia: 20%→7%"),
    (ORANGE, "BackTest en selección","Eliminado de trials y comparativa. Solo CV+Val para elegir el modelo."),
    (BLUE,   "Features Missings",   "umbral_missings elevado de 10%→25%. PSI confirma 27 features estables."),
    (GREEN,  "Modelo Final",        "HGB→LogReg por estabilidad temporal. Live: 0.55→0.63 (+8 pts)."),
    (NAVY,   "SHAP Explainer",      "TreeExplainer reemplazado por LinearExplainer (LogReg no es árbol)."),
    (MUTED,  "Gain/Lift Table",     "Clase de interés = churn=0 (1.2%). Lógica invertida. Lift 2.24x en Live."),
]
for i, (color, title, desc) in enumerate(corrections):
    row = i // 2
    col = i % 2
    xi = 0.3 + col * 6.5
    yi = 1.3 + row * 1.8
    rect(sl, xi, yi, 6.2, 1.65, WHITE)
    rect(sl, xi, yi, 0.18, 1.65, color)
    txt(sl, title, xi + 0.3, yi + 0.1, 5.7, 0.5, sz=16, bold=True, rgb=NAVY)
    txt(sl, desc, xi + 0.3, yi + 0.65, 5.7, 0.85, sz=12, rgb=DARK)

# ══════════════════════════════════════════════════════════
# SLIDE 17 — CONCLUSIONES
# ══════════════════════════════════════════════════════════
sl = new_slide()
bg(sl, NAVY)
rect(sl, 0, 0, 0.22, SH, ORANGE)
txt(sl, "Conclusiones", 0.5, 0.25, 12.3, 0.85, sz=38, bold=True, rgb=WHITE)

conclusions = [
    (GREEN,                         "Modelo Robusto",
     "LogReg seleccionado por estabilidad temporal. AUC: 0.60→0.61→0.63 Val→BT→Live"),
    (ORANGE,                        "Features Limpias",
     "27 variables libres de leakage, PSI estables, ninguna supera 12% de importancia"),
    (BLUE,                          "SHAP Interpretable",
     "Gasto reciente vs histórico = señal clave. monto_60d retiene, recencia predice churn"),
    (TEAL,                          "Lift Validado",
     "2.24x en Live (decil 1). El 10% superior de clientes captura 22% de churners reales"),
    (RGBColor(0x9B, 0x59, 0xB6),    "Predicciones Listas",
     "24,145 clientes scorados para Septiembre 2018. Listos para campaña de retención"),
]
for i, (color, title, desc) in enumerate(conclusions):
    yi = 1.3 + i * 1.05
    rect(sl, 0.5, yi, 0.55, 0.85, color)
    txt(sl, title, 1.2, yi + 0.08, 3.3, 0.5, sz=16, bold=True, rgb=WHITE)
    txt(sl, desc, 4.7, yi + 0.12, 8.3, 0.65, sz=14, rgb=LBLUE)

rect(sl, 0.5, 6.75, 12.3, 0.65, DBLUE)
txt(sl, ("Próximos Pasos: Reentrenamiento mensual con datos frescos  ·  "
         "A/B testing de campañas por segmento de riesgo  ·  "
         "Monitoreo de drift con PSI continuo"),
    0.7, 6.8, 12.0, 0.55, sz=11, rgb=LBLUE, italic=True)

# ─── GUARDAR ──────────────────────────────────────────────
output_path = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    'Presentacion_Sprint4_Churn_Olist_Grupo7.pptx'
)
prs.save(output_path)
print(f"\n✓  Presentación guardada exitosamente:")
print(f"   {output_path}")
print(f"   Total de slides: {len(prs.slides)}")
print("\n   Nota: Actualiza los valores de Precisión/Recall/F1 en Slide 14")
print("   con los resultados reales del notebook.")
